import tkinter as tk
from tkinter import filedialog, Menu, messagebox, StringVar, IntVar
import hashlib
import os
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import pandas as pd
from pptx import Presentation
from moviepy.editor import VideoFileClip
from cryptography.fernet import Fernet
import requests
from PyQt5.QtWidgets import QApplication, QMainWindow, QAction, QLineEdit, QLabel
from PyQt5.QtWebEngineWidgets import QWebEngineView
import pygame

import datetime
import json
import platform
import subprocess

# Define constants
NEWS_API_KEY = 'YOUR_NEWSAPI_KEY'
EMAIL_SERVER = 'smtp.example.com'
EMAIL_PORT = 587
EMAIL_USER = 'user@example.com'
EMAIL_PASS = 'password'
FERNET_KEY = Fernet.generate_key()
ANTIVIRUS_SIGNATURES = {
    "eicar_test": "44d88612fea8a8f36de82e1278abb02f",
    "malware_sample_1": "098f6bcd4621d373cade4e832627b4f6",
    "malware_sample_2": "ad0234829205b9033196ba818f7a872b",
    # Add additional signatures as required
}

def calculate_md5(file_path):
    """Calculate the MD5 hash of the given file."""
    try:
        with open(file_path, "rb") as file:
            file_content = file.read()
            md5_hash = hashlib.md5(file_content).hexdigest()
            return md5_hash
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return None

def scan_file(file_path):
    """Scan a single file for known malware signatures."""
    file_hash = calculate_md5(file_path)
    if file_hash:
        for name, signature in ANTIVIRUS_SIGNATURES.items():
            if file_hash == signature:
                return name
    return None

def scan_directory(directory):
    """Recursively scan a directory for malware and display results."""
    malware_found = False
    for root, dirs, files in os.walk(directory):
        for file_name in files:
            file_path = os.path.join(root, file_name)
            malware_name = scan_file(file_path)
            if malware_name:
                print(f"Malware detected: {malware_name} in {file_path}")
                malware_found = True
            else:
                print(f"Scanned: {file_path} - No threats found")
    if not malware_found:
        print("No malware detected.")

def send_email(sender_email, receiver_email, subject, body):
    """Send an email with the provided details."""
    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['To'] = receiver_email
    msg['Subject'] = subject
    msg.attach(MIMEText(body, 'plain'))
    try:
        with smtplib.SMTP(EMAIL_SERVER, EMAIL_PORT) as server:
            server.starttls()
            server.login(sender_email, EMAIL_PASS)
            server.sendmail(sender_email, receiver_email, msg.as_string())
            print("Email sent successfully!")
    except Exception as e:
        print(f"Failed to send email: {e}")

def play_audio(file_path):
    """Play an audio file."""
    pygame.init()
    pygame.mixer.init()
    pygame.mixer.music.load(file_path)
    pygame.mixer.music.play()

def video_to_audio(video_path, output_path):
    """Convert video to audio."""
    video_clip = VideoFileClip(video_path)
    audio_clip = video_clip.audio
    audio_clip.write_audiofile(output_path)
    print(f"Audio file saved to {output_path}")

def open_text_file(file_path):
    """Open and display the contents of a text file."""
    try:
        with open(file_path, 'r') as file:
            print(file.read())
    except Exception as e:
        print(f"Error opening file: {e}")

def open_pdf(file_path):
    """Open and display a PDF file."""
    try:
        pdf_document = fitz.open(file_path)
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text = page.get_text()
            print(f"Page {page_num+1}:")
            print(text)
    except Exception as e:
        print(f"Error opening PDF: {e}")

def open_spreadsheet(file_path):
    """Open and display the contents of a spreadsheet."""
    try:
        df = pd.read_excel(file_path)
        print(df)
    except Exception as e:
        print(f"Error opening spreadsheet: {e}")

def open_presentation(file_path):
    """Open and display the contents of a presentation."""
    try:
        presentation = Presentation(file_path)
        for slide_num, slide in enumerate(presentation.slides):
            print(f"Slide {slide_num+1}:")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    print(shape.text)
    except Exception as e:
        print(f"Error opening presentation: {e}")

def encrypt_file(file_path, key):
    """Encrypt a file using the provided key."""
    fernet = Fernet(key)
    try:
        with open(file_path, 'rb') as file:
            data = file.read()
        encrypted_data = fernet.encrypt(data)
        with open(file_path + ".enc", 'wb') as file:
            file.write(encrypted_data)
        print(f"File encrypted: {file_path}.enc")
    except Exception as e:
        print(f"Error encrypting file: {e}")

def decrypt_file(file_path, key):
    """Decrypt a file using the provided key."""
    fernet = Fernet(key)
    try:
        with open(file_path, 'rb') as file:
            encrypted_data = file.read()
        decrypted_data = fernet.decrypt(encrypted_data)
        with open(file_path.replace(".enc", ""), 'wb') as file:
            file.write(decrypted_data)
        print(f"File decrypted: {file_path.replace('.enc', '')}")
    except Exception as e:
        print(f"Error decrypting file: {e}")

def fetch_news():
    """Fetch and display news headlines."""
    try:
        response = requests.get(f"https://newsapi.org/v2/top-headlines?country=us&apiKey={NEWS_API_KEY}")
        news_data = response.json()
        for article in news_data['articles']:
            print(article['title'])
    except Exception as e:
        print(f"Error fetching news: {e}")

def setup_experience():
    """Setup experience for first-time users."""
    print("Welcome to SierraOE! Let's get started with the setup.")
    user_name = input("Enter your name: ")
    user_email = input("Enter your email: ")
    print(f"Setup complete! Welcome, {user_name}. Your system is now ready.")
    with open("setup_complete.txt", 'w') as f:
        f.write("Setup completed.")

# SierraOE GUI Implementation
class SierraOE(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("SierraOE")
        self.geometry("1024x768")
        self.configure(bg="#2e2e2e")
        self.create_menu()
        self.create_dock()
        self.create_desktop()

        # First-time setup
        if not os.path.exists("setup_complete.txt"):
            setup_experience()

    def create_menu(self):
        menu_bar = Menu(self)
        self.config(menu=menu_bar)
        
        file_menu = Menu(menu_bar, tearoff=0)
        menu_bar.add_cascade(label="File", menu=file_menu)
        file_menu.add_command(label="Exit", command=self.quit)
        
        utilities_menu = Menu(menu_bar, tearoff=0)
        menu_bar.add_cascade(label="Utilities", menu=utilities_menu)
        utilities_menu.add_command(label="Antivirus", command=self.open_antivirus)
        utilities_menu.add_command(label="Settings", command=self.open_settings)
        utilities_menu.add_command(label="Mini-Games", command=self.open_mini_games)
        utilities_menu.add_command(label="News and Updates", command=self.open_news_updates)
        
        apps_menu = Menu(menu_bar, tearoff=0)
        menu_bar.add_cascade(label="Applications", menu=apps_menu)
        apps_menu.add_command(label="Code Editors", command=self.open_code_editors)
        apps_menu.add_command(label="Office Suite", command=self.open_office_suite)
        apps_menu.add_command(label="Email Client", command=self.open_email_client)
        apps_menu.add_command(label="Sierra WebWorld", command=self.open_web_browser)

    def create_dock(self):
        dock_frame = tk.Frame(self, bg="#1c1c1c")
        dock_frame.pack(side="bottom", fill="x")
        
        tk.Button(dock_frame, text="File Explorer", command=self.open_file_explorer).pack(side="left", padx=5, pady=5)
        tk.Button(dock_frame, text="Media Player", command=self.open_media_player).pack(side="left", padx=5, pady=5)
        tk.Button(dock_frame, text="Settings", command=self.open_settings).pack(side="left", padx=5, pady=5)
        tk.Button(dock_frame, text="News", command=self.open_news_updates).pack(side="left", padx=5, pady=5)

    def create_desktop(self):
        desktop_frame = tk.Frame(self, bg="#2e2e2e")
        desktop_frame.pack(expand=True, fill="both")
        tk.Label(desktop_frame, text="Welcome to SierraOE", fg="white", bg="#2e2e2e", font=("Arial", 24)).pack()

    def open_file_explorer(self):
        file_path = filedialog.askopenfilename()
        if file_path:
            file_ext = os.path.splitext(file_path)[1]
            if file_ext == ".txt":
                open_text_file(file_path)
            elif file_ext == ".pdf":
                open_pdf(file_path)
            elif file_ext in [".xls", ".xlsx"]:
                open_spreadsheet(file_path)
            elif file_ext in [".pptx"]:
                open_presentation(file_path)
            else:
                print("Unsupported file type.")

    def open_media_player(self):
        file_path = filedialog.askopenfilename()
        if file_path:
            if file_path.endswith(".mp3"):
                play_audio(file_path)
            elif file_path.endswith(".mp4"):
                video_to_audio(file_path, file_path.replace(".mp4", ".mp3"))

    def open_email_client(self):
        print("Opening email client...")

    def open_code_editors(self):
        print("Opening code editors...")

    def open_office_suite(self):
        print("Opening office suite...")

    def open_settings(self):
        print("Opening settings...")

    def open_mini_games(self):
        print("Opening mini-games...")

    def open_news_updates(self):
        fetch_news()

    def open_web_browser(self):
        app = QApplication([])
        window = QMainWindow()
        browser = QWebEngineView()
        browser.setUrl("http://www.google.com")
        window.setCentralWidget(browser)
        window.show()
        app.exec_()

    def open_antivirus(self):
        directory = filedialog.askdirectory()
        if directory:
            scan_directory(directory)

if __name__ == "__main__":
    app = SierraOE()
    app.mainloop()
