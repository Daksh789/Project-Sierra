import tkinter as tk
from tkinter import filedialog, simpledialog, colorchooser, messagebox
import datetime
import webbrowser
import platform
import psutil
import requests
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
import csv
import os

# Define SierraOE Class
class SierraOE(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("SierraOE")
        self.geometry("1200x800")  # Initial size, will adjust to screen
        self.configure(bg="black")

        self.create_widgets()
        self.create_menu()

    def create_widgets(self):
        # Create a Text widget for output area
        self.output_area = tk.Text(self, height=20, width=100)
        self.output_area.pack(pady=10, padx=10, fill=tk.BOTH, expand=True)

    def create_menu(self):
        menu = tk.Menu(self)
        self.config(menu=menu)

        file_menu = tk.Menu(menu)
        menu.add_cascade(label="File", menu=file_menu)
        file_menu.add_command(label="Exit", command=self.quit)

        tools_menu = tk.Menu(menu)
        menu.add_cascade(label="Tools", menu=tools_menu)
        tools_menu.add_command(label="Python Editor", command=self.python_editor)
        tools_menu.add_command(label="C++ Editor", command=self.cpp_editor)
        tools_menu.add_command(label="JavaScript Editor", command=self.js_editor)
        tools_menu.add_command(label="Spreadsheet", command=self.spreadsheet)
        tools_menu.add_command(label="Word Processor", command=self.word_processor)
        tools_menu.add_command(label="Presentation Manager", command=self.presentation_manager)
        tools_menu.add_command(label="Video Editor", command=self.video_editor)
        tools_menu.add_command(label="Email", command=self.email_gui)
        tools_menu.add_command(label="Timer", command=self.timer)
        tools_menu.add_command(label="Alarm", command=self.alarm)
        tools_menu.add_command(label="Antivirus", command=self.antivirus)
        tools_menu.add_command(label="Device Manager", command=self.device_manager)
        tools_menu.add_command(label="Web Browser", command=self.open_web_browser)
        tools_menu.add_command(label="WhatsApp Web", command=self.open_whatsapp_web)
        tools_menu.add_command(label="Zoom", command=self.open_zoom)
        tools_menu.add_command(label="Khan Academy", command=self.open_khan_academy)
        tools_menu.add_command(label="Change Color", command=self.change_color)
        tools_menu.add_command(label="Show Time", command=self.show_time)
        tools_menu.add_command(label="Show Date", command=self.show_date)
        tools_menu.add_command(label="Get Weather", command=self.get_weather)
        tools_menu.add_command(label="System Details", command=self.system_details)
        tools_menu.add_command(label="Help", command=self.help_utility)
        tools_menu.add_command(label="GUI", command=self.show_gui)

    def show_gui(self):
        self.hide_gui = tk.Toplevel(self)
        self.hide_gui.title("SierraOE GUI")
        self.hide_gui.geometry("1200x800")

        # Adding functionality for GUI
        menu = tk.Menu(self.hide_gui)
        self.hide_gui.config(menu=menu)

        file_menu = tk.Menu(menu)
        menu.add_cascade(label="File", menu=file_menu)
        file_menu.add_command(label="Exit", command=self.hide_gui.destroy)

        tools_menu = tk.Menu(menu)
        menu.add_cascade(label="Tools", menu=tools_menu)
        tools_menu.add_command(label="Python Editor", command=self.python_editor)
        tools_menu.add_command(label="C++ Editor", command=self.cpp_editor)
        tools_menu.add_command(label="JavaScript Editor", command=self.js_editor)
        tools_menu.add_command(label="Spreadsheet", command=self.spreadsheet)
        tools_menu.add_command(label="Word Processor", command=self.word_processor)
        tools_menu.add_command(label="Presentation Manager", command=self.presentation_manager)
        tools_menu.add_command(label="Video Editor", command=self.video_editor)
        tools_menu.add_command(label="Email", command=self.email_gui)
        tools_menu.add_command(label="Timer", command=self.timer)
        tools_menu.add_command(label="Alarm", command=self.alarm)
        tools_menu.add_command(label="Antivirus", command=self.antivirus)
        tools_menu.add_command(label="Device Manager", command=self.device_manager)
        tools_menu.add_command(label="Web Browser", command=self.open_web_browser)
        tools_menu.add_command(label="WhatsApp Web", command=self.open_whatsapp_web)
        tools_menu.add_command(label="Zoom", command=self.open_zoom)
        tools_menu.add_command(label="Khan Academy", command=self.open_khan_academy)
        tools_menu.add_command(label="Change Color", command=self.change_color)
        tools_menu.add_command(label="Show Time", command=self.show_time)
        tools_menu.add_command(label="Show Date", command=self.show_date)
        tools_menu.add_command(label="Get Weather", command=self.get_weather)
        tools_menu.add_command(label="System Details", command=self.system_details)
        tools_menu.add_command(label="Help", command=self.help_utility)

        self.hide_gui.protocol("WM_DELETE_WINDOW", self.hide_gui.destroy)
        self.hide_gui.mainloop()

    # Textbox Functions
    def run_dos_command(self):
        command = simpledialog.askstring("Run DOS Command", "Enter DOS command:")
        if command:
            result = os.popen(command).read()
            self.output_area.insert(tk.END, f"Output:\n{result}\n")

    # Code Editors
    def python_editor(self):
        editor = tk.Toplevel(self)
        editor.title("Sierra Python Editor")
        editor.geometry("800x600")
        text_area = tk.Text(editor, wrap='word')
        text_area.pack(fill="both", expand=True)
        text_area.insert(tk.END, "# Write your Python code here...\n")

        # Add functionality to execute Python code
        def run_code():
            code = text_area.get("1.0", tk.END)
            try:
                exec(code)
                self.output_area.insert(tk.END, "Code executed successfully.\n")
            except Exception as e:
                self.output_area.insert(tk.END, f"Error: {e}\n")

        tk.Button(editor, text="Run Code", command=run_code).pack(pady=5)

    def cpp_editor(self):
        editor = tk.Toplevel(self)
        editor.title("Sierra C++ Editor")
        editor.geometry("800x600")
        text_area = tk.Text(editor, wrap='word')
        text_area.pack(fill="both", expand=True)
        text_area.insert(tk.END, "// Write your C++ code here...\n")

        # Add functionality to compile C++ code
        def compile_code():
            code = text_area.get("1.0", tk.END)
            with open("temp.cpp", "w") as file:
                file.write(code)
            os.system("g++ temp.cpp -o temp.exe")
            os.system("temp.exe")
            self.output_area.insert(tk.END, "Code compiled and executed successfully.\n")

        tk.Button(editor, text="Compile & Run", command=compile_code).pack(pady=5)

    def js_editor(self):
        editor = tk.Toplevel(self)
        editor.title("Sierra JS Editor")
        editor.geometry("800x600")
        text_area = tk.Text(editor, wrap='word')
        text_area.pack(fill="both", expand=True)
        text_area.insert(tk.END, "// Write your JavaScript code here...\n")

        # Add functionality to execute JavaScript code
        def run_code():
            code = text_area.get("1.0", tk.END)
            # For a real JS execution, you'd use a JS engine or web browser
            self.output_area.insert(tk.END, f"JavaScript code:\n{code}\n")

        tk.Button(editor, text="Run Code", command=run_code).pack(pady=5)

    # Office Suite
    def word_processor(self):
        word_window = tk.Toplevel(self)
        word_window.title("Sierra Word Processor")
        word_window.geometry("800x600")

        doc_area = tk.Text(word_window, height=30, width=100, wrap='word')
        doc_area.pack(pady=5)

        def save_document():
            file_path = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word Documents", "*.docx")])
            if file_path:
                from docx import Document
                doc = Document()
                doc.add_paragraph(doc_area.get("1.0", tk.END))
                doc.save(file_path)

        tk.Button(word_window, text="Save Document", command=save_document).pack(pady=5)

    def spreadsheet(self):
        spreadsheet_window = tk.Toplevel(self)
        spreadsheet_window.title("Sierra Spreadsheet")
        spreadsheet_window.geometry("800x600")

        table = tk.Text(spreadsheet_window, height=30, width=100, wrap='none')
        table.pack(pady=5)
        table.insert(tk.END, "Enter your spreadsheet data here...\n")

        def save_spreadsheet():
            file_path = filedialog.asksaveasfilename(defaultextension=".csv", filetypes=[("CSV Files", "*.csv")])
            if file_path:
                with open(file_path, "w", newline="") as file:
                    writer = csv.writer(file)
                    data = table.get("1.0", tk.END).splitlines()
                    for row in data:
                        writer.writerow(row.split())

        tk.Button(spreadsheet_window, text="Save Spreadsheet", command=save_spreadsheet).pack(pady=5)

    def presentation_manager(self):
        presentation_window = tk.Toplevel(self)
        presentation_window.title("Sierra Presentation Manager")
        presentation_window.geometry("800x600")
        text_area = tk.Text(presentation_window, wrap='word')
        text_area.pack(fill="both", expand=True)
        text_area.insert(tk.END, "Create your presentation slides here...\n")

        def save_presentation():
            file_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Presentations", "*.pptx")])
            if file_path:
                from pptx import Presentation
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                text_box = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
                text_frame = text_box.text_frame
                text_frame.text = text_area.get("1.0", tk.END)
                prs.save(file_path)

        tk.Button(presentation_window, text="Save Presentation", command=save_presentation).pack(pady=5)

    # Utilities
    def change_color(self):
        color_code = colorchooser.askcolor()[1]
        if color_code:
            self.configure(bg=color_code)

    def show_time(self):
        self.output_area.insert(tk.END, f"Current Time: {datetime.datetime.now().strftime('%H:%M:%S')}\n")

    def show_date(self):
        self.output_area.insert(tk.END, f"Current Date: {datetime.datetime.now().strftime('%Y-%m-%d')}\n")

    def get_weather(self):
        location = simpledialog.askstring("Weather", "Enter city:")
        if location:
            api_key = "13a8a9d380e9d4f74c8b9bca9420b196"
            response = requests.get(f"http://api.openweathermap.org/data/2.5/weather?q={location}&appid={api_key}&units=metric")
            data = response.json()
            if response.status_code == 200:
                weather = data['weather'][0]['description']
                temp = data['main']['temp']
                self.output_area.insert(tk.END, f"Weather in {location}: {weather}, {temp}°C\n")
            else:
                self.output_area.insert(tk.END, "Failed to get weather data.\n")

    def system_details(self):
        details = (
            f"System: {platform.system()}\n"
            f"Node Name: {platform.node()}\n"
            f"Release: {platform.release()}\n"
            f"Version: {platform.version()}\n"
            f"Machine: {platform.machine()}\n"
            f"Processor: {platform.processor()}\n"
            f"CPU Cores: {psutil.cpu_count(logical=False)}\n"
            f"RAM: {round(psutil.virtual_memory().total / (1024**3))} GB\n"
        )
        self.output_area.insert(tk.END, details)

    def help_utility(self):
        help_window = tk.Toplevel(self)
        help_window.title("Help Utility")
        help_window.geometry("400x300")
        help_text = tk.Text(help_window)
        help_text.pack(fill="both", expand=True)
        help_text.insert(tk.END, "Help content goes here...\n")

    def antivirus(self):
        antivirus_window = tk.Toplevel(self)
        antivirus_window.title("Sierra Antivirus")
        antivirus_window.geometry("800x600")
        text_area = tk.Text(antivirus_window)
        text_area.pack(fill="both", expand=True)
        text_area.insert(tk.END, "Antivirus tools and scanning logic...\n")

    def device_manager(self):
        device_manager_window = tk.Toplevel(self)
        device_manager_window.title("Sierra Device Manager")
        device_manager_window.geometry("800x600")
        text_area = tk.Text(device_manager_window)
        text_area.pack(fill="both", expand=True)
        text_area.insert(tk.END, "Device manager tools and management logic...\n")

    def presentation_manager(self):
        presentation_window = tk.Toplevel(self)
        presentation_window.title("Sierra Presentation Manager")
        presentation_window.geometry("800x600")
        text_area = tk.Text(presentation_window)
        text_area.pack(fill="both", expand=True)
        text_area.insert(tk.END, "Presentation management tools...\n")

    def video_editor(self):
        video_editor_window = tk.Toplevel(self)
        video_editor_window.title("Sierra Video Editor")
        video_editor_window.geometry("800x600")
        text_area = tk.Text(video_editor_window)
        text_area.pack(fill="both", expand=True)
        text_area.insert(tk.END, "Video editing tools and functionality...\n")

    def email_gui(self):
        email_window = tk.Toplevel(self)
        email_window.title("Sierra Email")
        email_window.geometry("600x400")

        tk.Label(email_window, text="To:").pack(pady=5)
        to_entry = tk.Entry(email_window, width=50)
        to_entry.pack(pady=5)

        tk.Label(email_window, text="Subject:").pack(pady=5)
        subject_entry = tk.Entry(email_window, width=50)
        subject_entry.pack(pady=5)

        tk.Label(email_window, text="Message:").pack(pady=5)
        message_text = tk.Text(email_window, height=10, width=50)
        message_text.pack(pady=5)

        def send_email():
            to_address = to_entry.get()
            subject = subject_entry.get()
            message = message_text.get("1.0", tk.END)
            from_address = "your_email@example.com"  # Replace with your email
            password = "your_password"  # Replace with your email password

            msg = MIMEMultipart()
            msg['From'] = from_address
            msg['To'] = to_address
            msg['Subject'] = subject
            msg.attach(MIMEText(message, 'plain'))

            try:
                server = smtplib.SMTP('smtp.example.com', 587)  # Replace with your SMTP server and port
                server.starttls()
                server.login(from_address, password)
                server.sendmail(from_address, to_address, msg.as_string())
                server.quit()
                self.output_area.insert(tk.END, "Email sent successfully.\n")
            except Exception as e:
                self.output_area.insert(tk.END, f"Failed to send email: {e}\n")

        tk.Button(email_window, text="Send Email", command=send_email).pack(pady=5)

    def open_web_browser(self):
        webbrowser.open("http://www.google.com")

    def open_whatsapp_web(self):
        webbrowser.open("https://web.whatsapp.com")

    def open_zoom(self):
        webbrowser.open("https://zoom.us")

    def open_khan_academy(self):
        webbrowser.open("https://www.khanacademy.org")

    def timer(self):
        timer_window = tk.Toplevel(self)
        timer_window.title("Sierra Timer")
        timer_window.geometry("400x300")

        tk.Label(timer_window, text="Timer (seconds):").pack(pady=5)
        timer_entry = tk.Entry(timer_window)
        timer_entry.pack(pady=5)

        def start_timer():
            try:
                seconds = int(timer_entry.get())
                if seconds < 0:
                    raise ValueError("Negative time not allowed.")
                self.output_area.insert(tk.END, f"Timer set for {seconds} seconds...\n")
                self.update()
                self.after(seconds * 1000, lambda: self.output_area.insert(tk.END, "Timer finished!\n"))
            except ValueError as e:
                self.output_area.insert(tk.END, f"Invalid input: {e}\n")

        tk.Button(timer_window, text="Start Timer", command=start_timer).pack(pady=5)

    def alarm(self):
        alarm_window = tk.Toplevel(self)
        alarm_window.title("Sierra Alarm")
        alarm_window.geometry("400x300")

        tk.Label(alarm_window, text="Set Alarm Time (HH:MM):").pack(pady=5)
        alarm_entry = tk.Entry(alarm_window)
        alarm_entry.pack(pady=5)

        def set_alarm():
            alarm_time = alarm_entry.get()
            try:
                alarm_time = datetime.datetime.strptime(alarm_time, "%H:%M").time()
                self.output_area.insert(tk.END, f"Alarm set for {alarm_time}...\n")
                while True:
                    now = datetime.datetime.now().time()
                    if now >= alarm_time:
                        self.output_area.insert(tk.END, "Alarm ringing!\n")
                        break
                    self.update()
            except ValueError:
                self.output_area.insert(tk.END, "Invalid time format. Please use HH:MM.\n")

        tk.Button(alarm_window, text="Set Alarm", command=set_alarm).pack(pady=5)

    def change_theme(self, theme):
        if theme == "dark":
            self.configure(bg="black")
            self.output_area.configure(bg="black", fg="white")
        elif theme == "light":
            self.configure(bg="white")
            self.output_area.configure(bg="white", fg="black")
        elif theme == "blue":
            self.configure(bg="lightblue")
            self.output_area.configure(bg="lightblue", fg="black")

# Run the application
if __name__ == "__main__":
    app = SierraOE()
    app.mainloop()
